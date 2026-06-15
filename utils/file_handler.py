"""
File Handler — 文件类型检测 + 文本提取 + 自动元数据提取

支持 8 种核心格式的分层处理：
  层级1 (自带元数据): EPUB / PDF(文字层) / HTML
  层级2 (有文本无元数据): TXT / MD / JSON / CSV / SRT / DOCX / PPTX
  层级3 (需OCR): 图片 / 扫描PDF
  层级4 (无文件): 手动输入

用法:
  from utils.file_handler import detect_file_type, extract_text, extract_auto_metadata, detect_encoding

  file_info = detect_file_type("book.epub")
  # → {"tier": 1, "format": "epub", "extraction": "ebooklib", "has_auto_metadata": True}

  text = extract_text("book.epub", file_info)
  metadata = extract_auto_metadata("book.epub", file_info)
  # → {"title": "齿轮设计手册", "author": "张三", "metadata_source": {...}}
"""

import os
import re
import hashlib

from typing import Optional

# ═══════════════════════════════════════
# 常量
# ═══════════════════════════════════════

SIZE_LIMIT_MB = 50
SIZE_LIMIT_BYTES = SIZE_LIMIT_MB * 1024 * 1024

# 文件类型注册表：扩展名 → {tier, extraction_method, mime, has_auto_metadata}
FILE_TYPE_REGISTRY = {
    # ── 层级1: 自带元数据的文件 ──
    ".epub": {"tier": 1, "extraction": "ebooklib", "mime": "application/epub+zip",
              "has_auto_metadata": True},
    ".html": {"tier": 1, "extraction": "bs4", "mime": "text/html",
              "has_auto_metadata": True},
    ".htm":  {"tier": 1, "extraction": "bs4", "mime": "text/html",
              "has_auto_metadata": True},
    # ── PDF: 双路径（pypdf 优先，失败转 OCR） ──
    ".pdf":  {"tier": "auto", "extraction": "pypdf_then_ocr", "mime": "application/pdf",
              "has_auto_metadata": True},
    # ── 层级2: 有文本但无元数据 ──
    ".txt":  {"tier": 2, "extraction": "direct_read", "mime": "text/plain",
              "has_auto_metadata": False},
    ".md":   {"tier": 2, "extraction": "direct_read", "mime": "text/markdown",
              "has_auto_metadata": False},
    ".json": {"tier": 2, "extraction": "direct_read", "mime": "application/json",
              "has_auto_metadata": False},
    ".csv":  {"tier": 2, "extraction": "direct_read", "mime": "text/csv",
              "has_auto_metadata": False},
    ".srt":  {"tier": 2, "extraction": "srt_parse", "mime": "text/plain",
              "has_auto_metadata": False},
    ".docx": {"tier": 2, "extraction": "python_docx", "mime":
              "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
              "has_auto_metadata": False},
    ".pptx": {"tier": 2, "extraction": "python_pptx", "mime":
              "application/vnd.openxmlformats-officedocument.presentationml.presentation",
              "has_auto_metadata": False},
    # ── 层级3: 图片（需 OCR） ──
    ".jpg":  {"tier": 3, "extraction": "ocr", "mime": "image/jpeg",
              "has_auto_metadata": False},
    ".jpeg": {"tier": 3, "extraction": "ocr", "mime": "image/jpeg",
              "has_auto_metadata": False},
    ".png":  {"tier": 3, "extraction": "ocr", "mime": "image/png",
              "has_auto_metadata": False},
    ".tiff": {"tier": 3, "extraction": "ocr", "mime": "image/tiff",
              "has_auto_metadata": False},
    ".bmp":  {"tier": 3, "extraction": "ocr", "mime": "image/bmp",
              "has_auto_metadata": False},
    ".webp": {"tier": 3, "extraction": "ocr", "mime": "image/webp",
              "has_auto_metadata": False},
}

# 图片文件头魔数验证
IMAGE_HEADERS = {
    b'\xff\xd8\xff': 'jpeg',
    b'\x89PNG\r\n\x1a\n': 'png',
    b'GIF8': 'gif',
    b'RIFF': 'webp',        # 需要进一步验证 WEBP 子类型
    b'BM': 'bmp',
    b'MM\x00*': 'tiff',     # Big-endian TIFF
    b'II*\x00': 'tiff',     # Little-endian TIFF
}
PDF_HEADER = b'%PDF'

# 显示名称映射
FORMAT_DISPLAY_NAMES = {
    "epub": "EPUB 电子书",
    "pdf": "PDF 文档",
    "html": "HTML 网页",
    "txt": "纯文本",
    "md": "Markdown",
    "json": "JSON",
    "csv": "CSV 表格",
    "srt": "SRT 字幕",
    "docx": "Word 文档",
    "pptx": "PowerPoint 幻灯片",
    "jpeg": "JPEG 图片",
    "png": "PNG 图片",
    "tiff": "TIFF 图片",
    "bmp": "BMP 图片",
    "webp": "WebP 图片",
}

TIER_NAMES = {
    1: "自带元数据",
    2: "纯文本（需AI标注）",
    3: "需OCR识别",
    4: "手动输入",
}


# ═══════════════════════════════════════
# 文件类型检测
# ═══════════════════════════════════════

def detect_file_type(file_path: str) -> dict:
    """
    检测文件类型，返回处理信息。

    Args:
        file_path: 文件路径

    Returns:
        {
            "ok": bool,
            "tier": 1|2|3|4|"auto"|None,
            "format": "epub"|"pdf"|...,
            "extraction": "ebooklib"|"direct_read"|...,
            "has_auto_metadata": bool,
            "mime": "application/pdf"|...,
            "display_name": "PDF 文档"|...,
            "tier_name": "自带元数据"|...,
            "is_supported": bool,
            "warning": str|None,       # 提示信息（如超大小）
            "error": str|None,         # 致命错误
        }
    """
    result = {
        "ok": True,
        "tier": None,
        "format": None,
        "extraction": None,
        "has_auto_metadata": False,
        "mime": None,
        "display_name": None,
        "tier_name": None,
        "is_supported": False,
        "warning": None,
        "error": None,
    }

    if not os.path.exists(file_path):
        result["ok"] = False
        result["error"] = f"文件不存在: {file_path}"
        return result

    ext = os.path.splitext(file_path)[1].lower()
    file_size = os.path.getsize(file_path)

    # 文件大小检查
    if file_size > SIZE_LIMIT_BYTES:
        result["warning"] = (
            f"文件较大（{file_size / 1024 / 1024:.1f} MB），"
            f"建议控制在 {SIZE_LIMIT_MB} MB 以内，处理可能较慢。"
        )

    # 查找注册表
    registry_entry = FILE_TYPE_REGISTRY.get(ext)
    if registry_entry is None:
        result["tier"] = 2
        result["format"] = ext.lstrip(".") if ext else "unknown"
        result["extraction"] = "direct_read_fallback"
        result["has_auto_metadata"] = False
        result["is_supported"] = False
        result["warning"] = (
            f"未识别的文件格式（{ext}），将尝试以纯文本方式读取。"
            f"支持的格式：epub, pdf, txt, md, srt, docx, pptx, html, jpg, png"
        )
        result["display_name"] = f"未知格式 ({ext})"
        result["tier_name"] = "未知"
        return result

    result.update({
        "tier": registry_entry["tier"],
        "format": ext.lstrip("."),
        "extraction": registry_entry["extraction"],
        "has_auto_metadata": registry_entry["has_auto_metadata"],
        "mime": registry_entry["mime"],
        "is_supported": True,
    })

    # 显示名称
    fmt = result["format"]
    result["display_name"] = FORMAT_DISPLAY_NAMES.get(fmt, fmt.upper())

    # 层级名称
    tier = result["tier"]
    if tier != "auto":
        result["tier_name"] = TIER_NAMES.get(tier, "未知")
    else:
        result["tier_name"] = "自动判断"

    # 内容验证：检查文件头是否与扩展名一致
    try:
        with open(file_path, "rb") as f:
            header = f.read(16)
    except Exception as e:
        result["ok"] = False
        result["error"] = f"无法读取文件: {e}"
        return result

    # PDF 头部验证
    if ext == ".pdf" and not header.startswith(PDF_HEADER):
        result["warning"] = (result["warning"] or "") + " 文件头不是有效的 PDF 签名，可能损坏。"
        # 不阻止，pypdf 会进一步报错

    # 图片头部验证
    if ext in (".jpg", ".jpeg", ".png", ".tiff", ".bmp", ".webp"):
        matched = False
        for magic, img_type in IMAGE_HEADERS.items():
            if header.startswith(magic):
                matched = True
                if ext in (".jpg", ".jpeg") and img_type != "jpeg":
                    result["warning"] = (result["warning"] or "") + f" 扩展名与文件头不匹配。"
                if ext == ".png" and img_type != "png":
                    result["warning"] = (result["warning"] or "") + f" 扩展名与文件头不匹配。"
                break
        if not matched:
            result["warning"] = (result["warning"] or "") + " 文件头不是有效图片格式，OCR 可能失败。"

    return result


# ═══════════════════════════════════════
# 编码检测
# ═══════════════════════════════════════

def detect_encoding(file_path: str, sample_bytes: int = 4096) -> str:
    """
    自动检测文件编码，使用 UTF-8 → GBK → GB2312 → latin-1 兜底链。

    Args:
        file_path: 文件路径
        sample_bytes: 采样字节数（前 N 字节用于检测）

    Returns:
        编码名称，如 "utf-8", "gbk", "latin-1"
    """
    # 编码检测链（按优先级）
    encoding_chain = ["utf-8", "gbk", "gb2312", "latin-1"]

    try:
        with open(file_path, "rb") as f:
            raw = f.read(sample_bytes)
    except Exception:
        return "utf-8"  # 兜底

    # 尝试 chardet（可选依赖，未安装则跳过）
    try:
        import chardet
        detected = chardet.detect(raw)
        if detected and detected.get("encoding"):
            enc = detected["encoding"].lower()
            # 统一常见别名
            enc_map = {
                "gb2312": "gbk", "gb18030": "gbk",
                "iso-8859-1": "latin-1", "windows-1252": "latin-1",
            }
            enc = enc_map.get(enc, enc)
            # 验证检测结果是否真能解码
            try:
                raw.decode(enc)
                return enc
            except (UnicodeDecodeError, LookupError):
                pass  # 检测不准，走兜底链
    except ImportError:
        pass  # chardet 未安装

    # 兜底链：逐个尝试
    for enc in encoding_chain:
        try:
            raw.decode(enc)
            return enc
        except (UnicodeDecodeError, LookupError):
            continue

    return "latin-1"  # 最终兜底（永远不会抛异常）


def _read_text_with_fallback(file_path: str) -> str:
    """使用编码检测链读取整个文本文件。"""
    enc = detect_encoding(file_path)
    with open(file_path, "r", encoding=enc, errors="replace") as f:
        return f.read()


# ═══════════════════════════════════════
# 文本提取（按格式分发）
# ═══════════════════════════════════════

def extract_text(file_path: str, file_info: dict = None) -> dict:
    """
    从文件提取纯文本内容。

    Args:
        file_path: 文件路径
        file_info: detect_file_type() 的结果，可选（None 时自动检测）

    Returns:
        {"ok": bool, "text": str, "error": str|None, "extraction_method": str}
    """
    if file_info is None:
        file_info = detect_file_type(file_path)
        if not file_info["ok"]:
            return {"ok": False, "text": "", "error": file_info["error"],
                    "extraction_method": None}

    fmt = file_info["format"]

    # 映射格式 → 提取函数
    extractors = {
        "epub": _extract_epub,
        "pdf":  _extract_pdf,
        "html": _extract_html,
        "htm":  _extract_html,
        "txt":  _extract_text_file,
        "md":   _extract_text_file,
        "json": _extract_text_file,
        "csv":  _extract_text_file,
        "srt":  _extract_srt,
        "docx": _extract_docx,
        "pptx": _extract_pptx,
    }

    extractor = extractors.get(fmt)
    if extractor is None:
        # 图片格式：返回空文本（留给 OCR 页面处理）
        if file_info["tier"] == 3:
            return {
                "ok": True,
                "text": "",
                "error": None,
                "extraction_method": "ocr_required",
                "ocr_required": True,
            }
        # 未知格式：兜底尝试
        try:
            text = _read_text_with_fallback(file_path)
            return {"ok": True, "text": text, "error": None,
                    "extraction_method": "direct_read_fallback"}
        except Exception as e:
            return {"ok": False, "text": "", "error": str(e), "extraction_method": None}

    try:
        result = extractor(file_path)
        if isinstance(result, str):
            result = {"ok": True, "text": result, "error": None}
        result.setdefault("extraction_method", file_info["extraction"])
        result.setdefault("ok", True)
        result.setdefault("error", None)
        return result
    except Exception as e:
        return {"ok": False, "text": "", "error": str(e),
                "extraction_method": file_info["extraction"]}


# ── 各格式提取函数 ──

def _extract_text_file(file_path: str) -> dict:
    """纯文本格式：TXT / MD / JSON / CSV"""
    text = _read_text_with_fallback(file_path)
    return {"ok": True, "text": text, "error": None}


def _extract_srt(file_path: str) -> dict:
    """SRT 字幕：去掉序号和时间戳，保留纯文本"""
    text = _read_text_with_fallback(file_path)
    # 去掉序号行和时间戳行
    lines = text.split("\n")
    cleaned = []
    for line in lines:
        line = line.strip()
        # 跳过序号（纯数字）
        if line.isdigit():
            continue
        # 跳过时间戳行 "00:00:01,000 --> 00:00:04,000"
        if re.match(r"^\d{2}:\d{2}:\d{2}[,.]\d{3}\s*-->\s*\d{2}:\d{2}:\d{2}[,.]\d{3}$", line):
            continue
        # 跳过 WEBVTT 头
        if line == "WEBVTT":
            continue
        if line:
            cleaned.append(line)
    return {"ok": True, "text": "\n".join(cleaned), "error": None}


def _extract_docx(file_path: str) -> dict:
    """Word 文档（.docx）"""
    try:
        import docx
        doc = docx.Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paragraphs)
        return {"ok": True, "text": text, "error": None}
    except ImportError:
        return {"ok": False, "text": "", "error": "请安装 python-docx: pip install python-docx"}


def _extract_pptx(file_path: str) -> dict:
    """PowerPoint 幻灯片（.pptx）"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        t = paragraph.text.strip()
                        if t:
                            slide_lines.append(t)
            if slide_lines:
                slides_text.append(f"--- 幻灯片 {i} ---\n" + "\n".join(slide_lines))
        text = "\n\n".join(slides_text)
        return {"ok": True, "text": text, "error": None}
    except ImportError:
        return {"ok": False, "text": "", "error": "请安装 python-pptx: pip install python-pptx"}


def _extract_html(file_path: str) -> dict:
    """HTML 网页：提取标题 + meta + body 文本"""
    text = _read_text_with_fallback(file_path)
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(text, "html.parser")
        # 去掉 script/style
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        body_text = soup.get_text(separator="\n")
        # 清理多余空行
        lines = [line.strip() for line in body_text.split("\n") if line.strip()]
        return {"ok": True, "text": "\n".join(lines), "error": None}
    except ImportError:
        # 无 BeautifulSoup，返回原始 HTML 文本
        return {"ok": True, "text": text, "error": None,
                "warning": "未安装 BeautifulSoup，HTML 标签未清理"}


def _extract_epub(file_path: str) -> dict:
    """EPUB 电子书：提取文本内容"""
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
    except ImportError:
        missing = []
        try:
            import ebooklib  # noqa: F811
        except ImportError:
            missing.append("ebooklib")
        try:
            from bs4 import BeautifulSoup  # noqa: F811
        except ImportError:
            missing.append("beautifulsoup4")
        return {"ok": False, "text": "",
                "error": f"请安装缺少的库: pip install {' '.join(missing)}"}

    try:
        book = epub.read_epub(file_path)
    except Exception as e:
        return {"ok": False, "text": "", "error": f"EPUB 读取失败: {e}"}

    chapters = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        try:
            soup = BeautifulSoup(item.get_content(), "html.parser")
            chapter_text = soup.get_text(separator="\n")
            lines = [line.strip() for line in chapter_text.split("\n") if line.strip()]
            if lines:
                chapters.append("\n".join(lines))
        except Exception:
            continue

    text = "\n\n".join(chapters)
    if not text:
        return {"ok": False, "text": "",
                "error": "未能从 EPUB 中提取到文本内容"}
    return {"ok": True, "text": text, "error": None}


def _extract_pdf(file_path: str) -> dict:
    """
    PDF 双路径提取:
    1. 先用 pypdf 提取文字层
    2. 文字层不足 → 标记为需要 OCR
    """
    try:
        from pypdf import PdfReader
    except ImportError:
        # pypdf 可能叫 PyPDF2
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            return {"ok": False, "text": "",
                    "error": "请安装 pypdf: pip install pypdf"}

    try:
        reader = PdfReader(file_path)
    except Exception as e:
        return {"ok": False, "text": "", "error": f"PDF 文件无法打开: {e}"}

    # 提取所有页的文本
    pages_text = []
    for i, page in enumerate(reader.pages):
        try:
            pt = page.extract_text()
            if pt and pt.strip():
                pages_text.append(pt.strip())
        except Exception:
            continue

    text = "\n\n".join(pages_text)

    # 判断是否有足够的文字层
    # 规则：总字符数 < 100 且页数 > 1 → 可能是扫描版
    if len(text) < 100 and len(reader.pages) > 1:
        return {
            "ok": True,
            "text": text or "",
            "error": None,
            "warning": "此 PDF 文字层很少，可能是扫描版，建议使用 OCR 识别。",
            "ocr_recommended": True,
            "total_pages": len(reader.pages),
            "has_text_layer": False,
        }

    return {
        "ok": True,
        "text": text,
        "error": None,
        "total_pages": len(reader.pages),
        "has_text_layer": bool(text.strip()),
    }


# ═══════════════════════════════════════
# 自动元数据提取（层级1：EPUB/PDF/HTML）
# ═══════════════════════════════════════

def extract_auto_metadata(file_path: str, file_info: dict = None) -> dict:
    """
    从文件中提取自动可获取的元数据。

    仅对层级1文件有效（EPUB/PDF/HTML）。
    返回的每条字段附带 metadata_source 标记。

    Args:
        file_path: 文件路径
        file_info: detect_file_type() 的结果

    Returns:
        {
            "ok": bool,
            "metadata": {
                "title": {"value": "...", "source": "file"},
                "author": {"value": "...", "source": "file"},
                ...
            },
            "flat": {  # 扁平版本，方便直接 merge
                "title": "...",
                "author": "...",
            },
            "source_count": int,  # 成功提取的字段数
        }
    """
    if file_info is None:
        file_info = detect_file_type(file_path)

    if not file_info.get("has_auto_metadata"):
        return {"ok": True, "metadata": {}, "flat": {}, "source_count": 0,
                "note": "此文件类型不支持自动元数据提取"}

    fmt = file_info["format"]

    extractors = {
        "epub": _metadata_epub,
        "pdf":  _metadata_pdf,
        "html": _metadata_html,
        "htm":  _metadata_html,
    }

    extractor = extractors.get(fmt)
    if extractor is None:
        return {"ok": False, "metadata": {}, "flat": {}, "source_count": 0,
                "error": f"不支持的元数据提取格式: {fmt}"}

    try:
        raw_meta = extractor(file_path)
    except Exception as e:
        return {"ok": False, "metadata": {}, "flat": {}, "source_count": 0,
                "error": f"元数据提取失败: {e}"}

    # 包装为 source-tagged 格式
    tagged = {}
    flat = {}
    for key, value in raw_meta.items():
        if value and str(value).strip():
            tagged[key] = {"value": str(value).strip(), "source": "file"}
            flat[key] = str(value).strip()

    return {
        "ok": True,
        "metadata": tagged,
        "flat": flat,
        "source_count": len(tagged),
    }


def _metadata_epub(file_path: str) -> dict:
    """从 EPUB 提取 Dublin Core 元数据"""
    try:
        from ebooklib import epub
    except ImportError:
        return {}

    try:
        book = epub.read_epub(file_path)
    except Exception:
        return {}

    meta = {}
    # Dublin Core 字段映射
    dc_map = {
        "title":   ("title",),
        "creator": ("author",),
        "publisher": ("publisher",),
        "identifier": ("isbn",),
        "language": ("language",),
        "date":     ("date",),
        "description": ("description",),
    }

    for dc_key, target_keys in dc_map.items():
        values = book.get_metadata("DC", dc_key)
        if values:
            # values 是 [(value, attrs), ...] 的列表
            val = values[0][0] if isinstance(values[0], tuple) else str(values[0])
            for target_key in target_keys:
                if val:
                    meta[target_key] = val

    # ISBN 特殊处理：identifier 可能包含 "urn:isbn:" 前缀
    if meta.get("isbn"):
        meta["isbn"] = meta["isbn"].replace("urn:isbn:", "").strip()

    return meta


def _metadata_pdf(file_path: str) -> dict:
    """从 PDF 提取 Document Info 元数据"""
    try:
        from pypdf import PdfReader
    except ImportError:
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            return {}

    try:
        reader = PdfReader(file_path)
        info = reader.metadata
        if info is None:
            return {}
    except Exception:
        return {}

    meta = {}
    field_map = {
        "/Title":    "title",
        "/Author":   "author",
        "/Subject":  "subject",
        "/Creator":  "creator",
        "/Producer": "producer",
    }

    for pdf_key, meta_key in field_map.items():
        val = getattr(info, pdf_key.strip("/").lower(), None) or info.get(pdf_key, None)
        if val and str(val).strip():
            meta[meta_key] = str(val).strip()

    # 页数
    meta["page_count"] = str(len(reader.pages))

    return meta


def _metadata_html(file_path: str) -> dict:
    """从 HTML 提取 title + meta 标签"""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return {}

    text = _read_text_with_fallback(file_path)
    try:
        soup = BeautifulSoup(text, "html.parser")
    except Exception:
        return {}

    meta = {}

    # <title>
    title_tag = soup.find("title")
    if title_tag and title_tag.string:
        meta["title"] = title_tag.string.strip()

    # <meta name="description">
    desc = soup.find("meta", attrs={"name": "description"})
    if desc and desc.get("content"):
        meta["description"] = desc["content"].strip()

    # <meta name="keywords">
    kw = soup.find("meta", attrs={"name": "keywords"})
    if kw and kw.get("content"):
        meta["keywords"] = kw["content"].strip()

    # <meta name="author">
    author = soup.find("meta", attrs={"name": "author"})
    if author and author.get("content"):
        meta["author"] = author["content"].strip()

    return meta


# ═══════════════════════════════════════
# 合并元数据（file > llm > manual）
# ═══════════════════════════════════════

def merge_metadata(file_meta: dict, llm_meta: dict, manual_meta: dict = None) -> dict:
    """
    按优先级合并元数据：文件自带 > LLM 推断 > 手动默认值。

    每个字段返回 value + source + confidence 三元组。

    Args:
        file_meta: extract_auto_metadata() 返回的 flat 字典
        llm_meta:  auto_classify() 返回的分类结果
        manual_meta: 用户手动填写的值（可选）

    Returns:
        {"title": {"value": "...", "source": "file", "confidence": 1.0}, ...}
    """
    if manual_meta is None:
        manual_meta = {}

    # 需要合并的字段列表
    mergeable_fields = ["title", "author", "keywords", "description"]

    result = {}

    for field in mergeable_fields:
        if field in file_meta and file_meta[field]:
            result[field] = {"value": file_meta[field], "source": "file",
                             "confidence": 1.0}
        elif field in llm_meta and llm_meta[field]:
            result[field] = {"value": llm_meta[field], "source": "llm",
                             "confidence": 0.7}
        elif field in manual_meta and manual_meta[field]:
            result[field] = {"value": manual_meta[field], "source": "manual",
                             "confidence": 0.5}

    return result
